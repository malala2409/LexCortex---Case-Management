"""LexCortex – Präsentation erstellen (FAU Template)
   Feature/Code im Wechsel, Code-Highlighting, lange Code-Blöcke geteilt"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE = "PP-Master-FAU_RW_4zu3.pptx"
OUTPUT = "LexCortex_Praesentation.pptx"

# ── Farben ─────────────────────────────────────────────────
FAU_DARK    = RGBColor(0x00, 0x2F, 0x5D)
FAU_BLUE    = RGBColor(0x00, 0x4E, 0x8C)
HIGHLIGHT   = RGBColor(0x03, 0x63, 0xA5)  # leuchtendes Blau für Key-Code
CODE_KEY    = RGBColor(0x0E, 0x4D, 0x8C)  # wichtige Code-Zeilen
CODE_NORMAL = RGBColor(0x47, 0x53, 0x69)  # normaler Code (gedimmt)
CODE_COMMENT= RGBColor(0x94, 0xA3, 0xB8)  # Kommentare
CODE_BG     = RGBColor(0xF1, 0xF5, 0xF9)
SCREEN_BG   = RGBColor(0xE2, 0xE8, 0xF0)

def add_screenshot_slide(slide, image_path, caption=None):
    """Screenshot aus Datei einfügen"""
    if os.path.exists(image_path):
        # Bildgröße berechnen (4:3 Slide = 10x7.5, Bildbereich ~9.2x5.5)
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.3), Inches(9.0), Inches(5.8))
    else:
        # Fallback: Platzhalter
        frame = slide.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.2))
        frame.fill.solid(); frame.fill.fore_color.rgb = SCREEN_BG
        frame.line.color.rgb = CODE_NORMAL; frame.line.width = Pt(1)
        hint = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(0.8))
        p = hint.text_frame.paragraphs[0]
        p.text = f"📸 {image_path} nicht gefunden"
        p.font.size = Pt(16); p.font.color.rgb = MUTED; p.alignment = PP_ALIGN.CENTER
DARK_TEXT    = RGBColor(0x1E, 0x29, 0x3B)
MUTED       = RGBColor(0x64, 0x74, 0x8B)
PAGE_NUM    = RGBColor(0x94, 0xA3, 0xB8)

prs = Presentation(TEMPLATE)
slide_layouts = prs.slide_layouts

# ── Layouts ────────────────────────────────────────────────
for i, layout in enumerate(slide_layouts):
    name = layout.name.lower()
    if 'titelbild' in name:       title_layout = layout
    if name == 'inhaltsfolie':     content_layout = layout
    if name == 'textfolie':        text_layout = layout
    if 'nur titel' in name and 'weiß' in name: blank_layout = layout

# Fallbacks
if 'title_layout' not in dir(): title_layout = slide_layouts[0]
if 'content_layout' not in dir(): content_layout = slide_layouts[13] if len(slide_layouts) > 13 else slide_layouts[4]
if 'text_layout' not in dir(): text_layout = slide_layouts[4]
if 'blank_layout' not in dir(): blank_layout = slide_layouts[17]

print(f"Layouts: Title={title_layout.name}, Content={content_layout.name}, Blank={blank_layout.name}")

# ── Alte Slides löschen ────────────────────────────────────
def remove_all_slides(pres):
    while len(pres.slides) > 0:
        rId = pres.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        pres.part.drop_rel(rId)
        pres.slides._sldIdLst.remove(pres.slides._sldIdLst[0])

remove_all_slides(prs)

# ── Hilfsfunktionen ────────────────────────────────────────
def add_feature_slide(slide, cards, x=Inches(0.6), y=Inches(1.3)):
    """
    Feature-Folie: Alle Inhalte in EINER Textbox.
    cards = [(titel, [(pt, pd), ...]), ...]
    """
    tb = slide.shapes.add_textbox(x, y, Inches(8.8), Inches(5.8))
    tf = tb.text_frame; tf.word_wrap = True
    clear_first = True
    for card_title, points in cards:
        # Abstand zwischen Karten (außer vor der ersten)
        if not clear_first:
            spacer = tf.add_paragraph()
            spacer.text = ""; spacer.font.size = Pt(10); spacer.space_after = Pt(0)

        # Karten-Titel (ersten leeren Absatz überschreiben, sonst neuen anlegen)
        if clear_first:
            p = tf.paragraphs[0]
            clear_first = False
        else:
            p = tf.add_paragraph()
        p.text = card_title
        p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = FAU_DARK
        p.space_after = Pt(6)

        # Bullet Points
        for pt, pd in points:
            p = tf.add_paragraph()
            p.text = f"▸ {pt}: {pd}"
            p.font.size = Pt(16); p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(5)

def add_code_slide(slide, section_title, code_items, x=Inches(0.5), y=Inches(1.35), page_info=None):
    """
    Code-Folie: Alles in EINER Textbox (Titel + Code).
    code_items: Liste von (text, typ):
      'H' = Highlight (blau+fett)  'K' = Key (dunkel+fett)
      'N' = Normal (grau)  'C' = Comment (hellgrau)  'E' = Empty
    """
    title_text = section_title
    if page_info:
        title_text += f"  [{page_info}]"

    n = len(code_items)
    line_h = 0.23
    code_h = min(n * line_h + 0.6, 5.5)

    # Erst Hintergrund-Rechteck (wird von Textbox überlagert)
    box_y = y + Inches(0.45)
    box = slide.shapes.add_shape(1, x, box_y, Inches(9.0), Inches(code_h))
    box.fill.solid(); box.fill.fore_color.rgb = CODE_BG; box.line.fill.background()

    # Dann Textbox (liegt über dem Hintergrund)
    tb = slide.shapes.add_textbox(x + Inches(0.08), y, Inches(8.84), Inches(0.5 + code_h))
    tf = tb.text_frame; tf.word_wrap = True

    # Code-Titel
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = FAU_DARK
    p.space_after = Pt(8)

    type_colors = {
        'H': CODE_KEY, 'K': DARK_TEXT, 'N': CODE_NORMAL,
        'C': CODE_COMMENT, 'E': CODE_COMMENT,
    }

    for text, typ in code_items:
        p = tf.add_paragraph()
        p.text = text if text else " "
        p.font.size = Pt(11)
        p.font.name = "Courier New"
        p.font.color.rgb = type_colors.get(typ, CODE_NORMAL)
        p.font.bold = (typ == 'H' or typ == 'K')
        p.space_after = Pt(1)

# ──────────────────────────────────────────────────────────
# FOLIE 1: Titel
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(title_layout)
for shape in s.shapes:
    if shape.is_placeholder:
        if 'title' in shape.name.lower() or 'titel' in shape.name.lower():
            shape.text = "LexCortex"
        elif 'subtitle' in shape.name.lower() or 'untertitel' in shape.name.lower():
            shape.text = "Fristen- und Fallmanagement für Anwälte"
        elif shape.has_text_frame and not shape.text.strip():
            shape.text = "Welche Funktionen wurden umgesetzt?"

# ──────────────────────────────────────────────────────────
# FOLIE 2: Agenda
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(content_layout)
if s.shapes.title: s.shapes.title.text = "Agenda"
body = None
for ph in s.placeholders:
    if ph.placeholder_format.idx != 0 and ph.has_text_frame:
        body = ph; break
if body is None:
    for sh in s.shapes:
        if sh.has_text_frame and sh != s.shapes.title: body = sh; break
if body:
    tf = body.text_frame; tf.clear()
    items = ["Dashboard – Übersicht & Kennzahlen",
             "Fallverwaltung & Phasen-Timeline",
             "Automatische Fristberechnung",
             "Kalender & Tagesübersicht",
             "Aufgaben – Inline-Bearbeitung",
             "JavaScript & jQuery im Einsatz"]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(16); p.font.color.rgb = DARK_TEXT; p.space_after = Pt(10)

# ──────────────────────────────────────────────────────────
# FOLIE 3: Dashboard (Feature)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Dashboard – Übersicht & Kennzahlen"
add_feature_slide(s, [
    ("📊 Vier Metriken auf der Startseite", [
        ("Aktive Fälle", "Anzahl der laufenden Verfahren — direkt aus der Datenbank abgefragt"),
        ("Fristen diese Woche", "Deadlines der nächsten 7 Tage im Überblick"),
        ("Überfällige Fristen", "Verpasste Deadlines werden rot hervorgehoben"),
        ("Abgeschlossene Fälle", "Erledigte Verfahren im aktuellen Quartal"),
    ]),
    ("📋 Cases-Tabelle mit Filterfunktion", [
        ("Phasen-Dropdown", "Filtert die Tabelle nach der aktiven Verfahrensphase"),
        ("Farbige Status-Punkte", "Rot = überfällig, Orange = ≤ 7 Tage, Grün = in Ordnung"),
        ("Frist-Badges", "Jedem Fall wird der aktuelle Fristtyp angezeigt (Notfrist, Gerichtsfrist, etc.)"),
    ]),
])

# ──────────────────────────────────────────────────────────
# FOLIE 4: Dashboard Screenshot
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "📸 Dashboard – Live-Ansicht"
add_screenshot_slide(s, "screenshots/01-dashboard.png")

# ──────────────────────────────────────────────────────────
# FOLIE 5: Dashboard Code (Teil 1: Metriken)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Dashboard – PHP/MySQL Metriken"
add_code_slide(s, "SQL-Abfragen für die 4 Kennzahlen", [
    ("<?php", 'N'),
    ("// ═══ Metriken: 4 SQL-Abfragen ═══════════════", 'C'),
    ("", 'E'),
    ('$aktive = $db->query(', 'K'),
    ('    "SELECT COUNT(*) FROM cases', 'H'),
    ('    WHERE status = \'aktiv\'")', 'H'),
    ('    ->fetchColumn();', 'K'),
    ("", 'E'),
    ('$diese_woche = $db->query(', 'K'),
    ('    "SELECT COUNT(*) FROM deadlines', 'H'),
    ('    WHERE deadline_date BETWEEN CURDATE()', 'H'),
    ('    AND DATE_ADD(CURDATE(), INTERVAL 7 DAY)', 'H'),
    ('    AND erledigt = 0")', 'H'),
    ('    ->fetchColumn();', 'K'),
    ("", 'E'),
    ('$ueberfaellig = $db->query(', 'K'),
    ('    "SELECT COUNT(*) FROM deadlines', 'H'),
    ('    WHERE deadline_date < CURDATE()', 'H'),
    ('    AND erledigt = 0")', 'H'),
    ('    ->fetchColumn();', 'K'),
    ("", 'E'),
    ('$abgeschl = $db->query(', 'K'),
    ('    "SELECT COUNT(*) FROM cases', 'N'),
    ('    WHERE status = \'abgeschlossen\'")', 'N'),
    ('    ->fetchColumn();', 'K'),
], page_info="1/2")

# ──────────────────────────────────────────────────────────
# FOLIE 5: Dashboard Code (Teil 2: Cases + Badges)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Dashboard – Cases & Frist-Badges"
add_code_slide(s, "JOIN + Fristtyp-Mapping", [
    ("// ═══ Cases mit aktiver Phase (JOIN) ═════════", 'C'),
    ('$cases = $db->query("', 'K'),
    ('    SELECT c.*, p.title AS phase_title,', 'H'),
    ('           p.frist_type, p.status', 'H'),
    ('    FROM cases c', 'H'),
    ('    LEFT JOIN phases p', 'H'),
    ('        ON p.case_id = c.id', 'H'),
    ('        AND p.status = \'active\'', 'H'),
    ('    WHERE c.status = \'aktiv\'', 'H'),
    ('    ORDER BY c.created_at DESC', 'H'),
    ('")->fetchAll();', 'K'),
    ("", 'E'),
    ("// ═══ Frist-Badge Mapping ═══════════════════", 'C'),
    ('function fristBadge($fristType, $status) {', 'K'),
    ('    $map = [', 'N'),
    ("        'notfrist'      => ['badge-notfrist',     'Notfrist'],", 'H'),
    ("        'gerichtsfrist' => ['badge-gerichtsfrist','Gerichtsfrist'],", 'H'),
    ("        'termin'        => ['badge-termin',       'Termin'],", 'H'),
    ("        'urteil'        => ['badge-urteil',       'Urteil'],", 'H'),
    ('    ];', 'N'),
    ('    return "<span class=\'badge {$cls}\'>$label</span>";', 'K'),
    ('}', 'K'),
], page_info="2/2")

# ──────────────────────────────────────────────────────────
# FOLIE 6: Phasen-Timeline (Feature)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Fallverwaltung & Phasen-Timeline"
add_feature_slide(s, [
    ("⚖️ Zehn Verfahrensphasen eines Zivilprozesses", [
        ("Vorkonfigurierte Reihenfolge", "Mahnung → Klageerhebung → Zustellung → Verteidigungsanzeige → Klageerwiderung → Replik → Güteverhandlung → Mündliche Verhandlung → Urteil → Rechtskraft"),
        ("Fortschrittsbalken", "Zeigt an, wie viele der zehn Phasen bereits abgeschlossen sind"),
        ("Phase abschließen", "Nach Eingabe eines Datums wird die nächste Phase automatisch aktiviert"),
    ]),
    ("📝 Weitere Funktionen der Timeline", [
        ("Phase zurücknehmen", "Letzte abgeschlossene Phase kann per Klick rückgängig gemacht werden"),
        ("Datum bearbeiten", "Phasen-Datum direkt in der Timeline per Klick änderbar"),
        ("jQuery slideToggle", "Auf- und Zuklappen der Phasen mit Animation"),
    ]),
])

# ──────────────────────────────────────────────────────────
# Screenshot: Phasen-Timeline
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "📸 Phasen-Timeline – Live-Ansicht"
add_screenshot_slide(s, "screenshots/03-timeline.png")

# ──────────────────────────────────────────────────────────
# FOLIE: Phasen-Timeline (Code)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Phase abschließen – AJAX & PHP"
add_code_slide(s, "complete_phase.php (vereinfacht)", [
    ("// ═══ AJAX: Phase abschließen ═══════════════════", 'C'),
    ("$phaseTitle = $_POST['phase_title'] ?? '';", 'N'),
    ("", 'E'),
    ("// 1. Phase als done markieren", 'C'),
    ('$db->prepare("UPDATE phases', 'K'),
    ("    SET status = 'done', phase_date = ?,", 'H'),
    ("    completed_at = NOW() WHERE id = ?\")", 'H'),
    ('    ->execute([$date, $phaseId]);', 'K'),
    ("", 'E'),
    ("// 2. Nächste pending → active", 'C'),
    ('$next = $db->prepare("SELECT id FROM phases', 'H'),
    ("    WHERE case_id = ? AND status = 'pending'", 'H'),
    ("    ORDER BY id ASC LIMIT 1\")->execute([$caseId]);", 'H'),
    ('$nextPhase = $next->fetch();', 'N'),
    ('if ($nextPhase) {', 'K'),
    ('    $db->prepare("UPDATE phases SET status = \'active\'', 'H'),
    ('        WHERE id = ?")->execute([$nextPhase["id"]]);', 'H'),
    ('}', 'K'),
    ("", 'E'),
    ("// 3. Automatische Fristberechnung", 'C'),
    ('if ($phaseTitle === "Zustellung") {', 'H'),
    ('    $frist = (new DateTime($date))', 'H'),
    ("        ->modify('+14 days')->format('Y-m-d');", 'H'),
    ('    // Deadline + Phase-Datum setzen', 'C'),
    ('    $db->prepare("INSERT INTO deadlines ...")->execute(...);', 'H'),
    ('}', 'K'),
])

# ──────────────────────────────────────────────────────────
# FOLIE 8: Automatische Fristberechnung (Feature)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Automatische Fristberechnung"
add_feature_slide(s, [
    ("⚡ Zustellungsdatum → +14 Tage Notfrist", [
        ("§ 276 ZPO", "Nach Zustellung der Klage muss der Beklagte innerhalb von 14 Tagen seine Verteidigungsbereitschaft anzeigen"),
        ("Automatische Berechnung", "Wird beim Anlegen eines Falls ein Zustellungsdatum eingegeben, berechnet das System das Fristdatum und legt einen entsprechenden Deadline-Eintrag an"),
        ("Live-Vorschau im Formular", "JavaScript-Funktion calcDeadlinePreview() zeigt das berechnete Datum sofort an"),
    ]),
    ("🏛️ Urteil → +30 Tage Berufungsfrist", [
        ("Automatische Deadline", "Wird die Phase \"Urteil\" abgeschlossen, wird automatisch eine Berufungsfrist von 30 Tagen eingetragen"),
        ("Fristtyp", "In beiden Fällen wird die Frist als \"Notfrist\" klassifiziert und im Kalender rot markiert"),
    ]),
])

# ──────────────────────────────────────────────────────────
# Screenshot: Fristberechnung
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "📸 Fall anlegen – Live-Ansicht"
add_screenshot_slide(s, "screenshots/02-neuer-fall.png")

# ──────────────────────────────────────────────────────────
# FOLIE: Fristberechnung Code (Teil 1: PHP)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Fristberechnung – PHP Backend"
add_code_slide(s, "add_case.php: Phasen + Auto-Frist", [
    ("// ═══ 10 Phasen der Zahlungsklage ═══════════════", 'C'),
    ('$phasen = [', 'N'),
    ("    ['Vorgerichtliche Mahnung',    null,            null],", 'N'),
    ("    ['Klageerhebung',              null,            null],", 'N'),
    ("    ['Zustellung',                 \$zustellungDatum,'schluessel'],", 'H'),
    ("    ['Verteidigungsanzeige',       null,            'notfrist'],", 'H'),
    ("    ['Klageerwiderung',            null,            'gerichtsfrist'],", 'N'),
    ("    // ... Replik, Güteverhandlung, Mündl. Verh.,", 'C'),
    ("    ['Urteil',                     null,            'urteil'],", 'H'),
    ("    ['Rechtskraft / Vollstreckung',null,            null],", 'N'),
    ("];", 'N'),
    ("", 'E'),
    ("// ═══ Auto-Frist: Zustellung + 14 Tage ═══════", 'C'),
    ('if ($zustellungDatum) {', 'K'),
    ('    $fristDatum = (new DateTime($zustellungDatum))', 'H'),
    ("        ->modify('+14 days')->format('Y-m-d');", 'H'),
    ('    $db->prepare("INSERT INTO deadlines', 'K'),
    ("        (case_id, title, deadline_date,", 'N'),
    ("         frist_type, auto_calculated)", 'N'),
    ("        VALUES (?, 'Verteidigungsanzeige', ?,", 'H'),
    ("        'notfrist', true)\")->execute([...]);", 'H'),
    ('}', 'K'),
    ("", 'E'),
    ("// Zusätzlich: Urteil → +30 Tage (in complete_phase.php)", 'C'),
], page_info="1/2")

# ──────────────────────────────────────────────────────────
# FOLIE 10: Fristberechnung Code (Teil 2: JavaScript)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Fristberechnung – JavaScript Vorschau"
add_code_slide(s, "Live-Datumsvorschau im Formular", [
    ("// ═══ Echtzeit-Fristvorschau (Vanilla JS) ═══════", 'C'),
    ("// Wird bei Eingabe des Zustellungsdatums getriggert", 'C'),
    ("function calcDeadlinePreview() {", 'K'),
    ("    var val = document.getElementById(", 'N'),
    ("        'f-zustellung').value;", 'N'),
    ("    if (!val) {", 'K'),
    ("        hint.textContent = '';", 'N'),
    ("        return;", 'N'),
    ("    }", 'K'),
    ("", 'E'),
    ("    var d = new Date(val);", 'H'),
    ("    d.setDate(d.getDate() + 14);  // ← +14 Tage", 'H'),
    ("", 'E'),
    ("    var fmt = d.toLocaleDateString('de-DE', {", 'N'),
    ("        day: '2-digit', month: '2-digit',", 'N'),
    ("        year: 'numeric'", 'N'),
    ("    });", 'N'),
    ("", 'E'),
    ("    // Live-Feedback im Formular anzeigen", 'C'),
    ("    hint.textContent =", 'H'),
    ("        '⚡ Verteidigungsanzeige fällig am: '", 'H'),
    ("        + fmt + ' (automatisch +14 Tage)';", 'H'),
    ("}", 'K'),
    ("", 'E'),
    ("// HTML: <input oninput='calcDeadlinePreview()'>", 'C'),
    ("// <div id='hint-zustellung'></div>", 'C'),
], page_info="2/2")

# ──────────────────────────────────────────────────────────
# FOLIE 11: Kalender (Feature)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Kalender – Monatsübersicht & Navigation"
add_feature_slide(s, [
    ("📅 Selbst programmierter Monatskalender (Vanilla JS)", [
        ("Monatsnavigation", "Vor- und Zurück-Buttons zum Wechseln des Monats"),
        ("Farbliche Kennzeichnung", "Rot = Notfrist, Orange = Gerichtsfrist, Blau = Termin, Lila = Urteil"),
        ("Task-Darstellung", "Aufgaben erscheinen mit Warnsymbolen: ⚡ bei ≤ 3 Tagen, ⚠ bei Überfälligkeit"),
    ]),
    ("🔍 Tagesübersicht und Verlinkung", [
        ("Tag-Klick", "Öffnet ein Modal mit allen Deadlines und Tasks des angeklickten Tages"),
        ("Kalender → Case-Detail", "Klick auf einen Task führt zur Fall-Detailseite, die zugehörige Phase wird automatisch aufgeklappt und der Task blinkt zur Hervorhebung"),
        ("Task-Erstellung aus dem Kalender", "Über den +-Button kann direkt aus dem Kalender eine neue Aufgabe für einen beliebigen Fall angelegt werden"),
    ]),
])

# ──────────────────────────────────────────────────────────
# Screenshot: Kalender
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "📸 Kalender – Live-Ansicht"
add_screenshot_slide(s, "screenshots/04-kalender.png")

# ──────────────────────────────────────────────────────────
# FOLIE: Kalender Code (Teil 1: Struktur)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Kalender – JavaScript (Struktur)"
add_code_slide(s, "renderCalendar() – Zellen aufbauen", [
    ("// ═══ Monatskalender (Vanilla JS) ═══════════════", 'C'),
    ("function renderCalendar() {", 'K'),
    ("    var firstDay = new Date(calYear, calMonth, 1);", 'N'),
    ("    var daysInMonth = new Date(calYear, calMonth + 1, 0)", 'N'),
    ("        .getDate();", 'N'),
    ("    var startDow = firstDay.getDay() || 7;", 'N'),
    ("    startDow--;  // Montag = 0", 'N'),
    ("", 'E'),
    ("    for (var i = 0; i < totalCells; i++) {", 'H'),
    ("        var cell = document.createElement('div');", 'H'),
    ("        cell.className = 'cal-day';", 'H'),
    ("", 'E'),
    ("        // Tagesnummer", 'C'),
    ("        var numDiv = document.createElement('div');", 'H'),
    ("        numDiv.className = 'cal-day-num';", 'H'),
    ("        numDiv.textContent = day;", 'H'),
    ("        cell.appendChild(numDiv);", 'H'),
    ("", 'E'),
    ("        // Heute markieren", 'C'),
    ("        var isToday = (day === today.getDate()", 'N'),
    ("            && m === today.getMonth()", 'N'),
    ("            && y === today.getFullYear());", 'N'),
    ("        if (isToday) cell.classList.add('today');", 'N'),
    ("        // ... Events + Click-Handler (siehe Teil 2)", 'C'),
    ("    }", 'H'),
    ("}", 'K'),
], page_info="1/2")

# ──────────────────────────────────────────────────────────
# FOLIE 13: Kalender Code (Teil 2: Events + Navigation)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Kalender – Events & Kalender→Task"
add_code_slide(s, "Event-Pills, Tagesübersicht, Task-Navigation", [
    ("// ═══ Events aus PHP (CAL_EVENTS als JSON) ═══════", 'C'),
    ("var dateStr = y + '-' + pad(m+1) + '-' + pad(day);", 'N'),
    ("var dayEvents = calEvents[dateStr] || [];", 'H'),
    ("dayEvents.forEach(function(ev) {", 'H'),
    ("    var pill = document.createElement('a');", 'H'),
    ("    pill.className = 'cal-event ' + ev.color;", 'H'),
    ("    pill.textContent = ev.label;", 'H'),
    ("    // Kalender → Case-Detail + Task-Highlight", 'C'),
    ("    pill.href = 'case_detail.php?id=' + ev.case_id", 'H'),
    ("        + '&task_id=' + ev.task_id + '&from=calendar';", 'H'),
    ("    cell.appendChild(pill);", 'H'),
    ("});", 'H'),
    ("", 'E'),
    ("// ═══ Tagesübersicht (Tag-Klick) ═══════════════", 'C'),
    ("cell.addEventListener('click', function(e) {", 'K'),
    ("    if (e.target.classList.contains('cal-event'))", 'N'),
    ("        return;  // Event-Pill = Navigation", 'N'),
    ("    openDayModal(dateStr, fmtDate);", 'H'),
    ("});", 'K'),
    ("", 'E'),
    ("// ═══ Kalender → Task-Navigation ═══════════════", 'C'),
    ("// In case_detail.php:", 'C'),
    ("function highlightTaskFromCalendar() {", 'K'),
    ("    // Phase per slideToggle() aufklappen", 'K'),
    ("    // taskEl.scrollIntoView({behavior:'smooth'})", 'H'),
    ("    // taskEl.classList.add('flash-highlight')", 'H'),
    ("    // → CSS-Animation: 2x blinken", 'H'),
    ("}", 'K'),
], page_info="2/2")

# ──────────────────────────────────────────────────────────
# FOLIE 14: Aufgaben (Feature)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Aufgaben – Inline-Bearbeitung"
add_feature_slide(s, [
    ("✅ Task-Verwaltung pro Verfahrensphase", [
        ("Task erstellen", "Zu jeder Phase können Aufgaben angelegt werden — über ein Inline-Formular oder aus dem Kalender heraus"),
        ("Erledigt-Status umschalten", "Checkbox zum An- und Abhaken — der Status wird per AJAX gespeichert"),
        ("Inline-Bearbeitung", "Beschreibung und Deadline können direkt im Task geändert werden, ohne die Seite zu verlassen"),
    ]),
    ("📋 Aufgaben auf dem Dashboard", [
        ("Heutige und überfällige Tasks", "Die sechs dringendsten Aufgaben werden auf der Startseite angezeigt"),
        ("Warnhinweise", "Tasks mit weniger als 3 Tagen erhalten ein orangenes Warnsymbol, überfällige Tasks ein rotes"),
        ("Live-Aktualisierung", "Neu angelegte Tasks erscheinen ohne Seiten-Reload sowohl im Dashboard als auch im Kalender"),
    ]),
])

# ──────────────────────────────────────────────────────────
# Screenshot: Tasks
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "📸 Aufgaben – Live-Ansicht"
add_screenshot_slide(s, "screenshots/05-tagesansicht.png")

# ──────────────────────────────────────────────────────────
# FOLIE: Aufgaben (Code)
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "Tasks – jQuery AJAX Toggle"
add_code_slide(s, "JS: completeTask() + PHP: Toggle-Logik", [
    ("// ═══ JavaScript: Task-Checkbox Toggle ═════════", 'C'),
    ("function completeTask(taskId, checkbox) {", 'K'),
    ("    $.ajax({", 'N'),
    ("        url:  'ajax/complete_task.php',", 'H'),
    ("        type: 'POST',", 'H'),
    ("        data: { task_id: taskId },", 'H'),
    ("        success: function(res) {", 'K'),
    ("            if (res.erledigt) {", 'H'),
    ("                $(checkbox).addClass('done').text('✓');", 'H'),
    ("                $(checkbox).closest('.phase-task-item')", 'H'),
    ("                    .addClass('task-done');", 'H'),
    ("            } else {", 'H'),
    ("                $(checkbox).removeClass('done').text('');", 'H'),
    ("                $(checkbox).closest('.phase-task-item')", 'H'),
    ("                    .removeClass('task-done');", 'H'),
    ("            }", 'K'),
    ("        }", 'N'),
    ("    });", 'N'),
    ("}", 'K'),
    ("", 'E'),
    ("// ═══ PHP: Toggle 1↔0 ═══════════════════════════", 'C'),
    ('$row = $db->prepare("SELECT erledigt FROM tasks', 'H'),
    ('    WHERE id = ?")->execute([$taskId])->fetch();', 'H'),
    ('$newStatus = $row["erledigt"] ? 0 : 1;  // Flip', 'H'),
    ('$db->prepare("UPDATE tasks SET erledigt = ?', 'H'),
    ('    WHERE id = ?")->execute([$newStatus, $taskId]);', 'H'),
    ('echo json_encode(["success" => true,', 'H'),
    ('    "erledigt" => (bool)$newStatus]);', 'H'),
])

# ──────────────────────────────────────────────────────────
# FOLIE 16: JavaScript & jQuery Übersicht
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
s.shapes.title.text = "JavaScript & jQuery – Interaktivität im Frontend"
add_feature_slide(s, [
    ("📜 Vanilla JavaScript", [
        ("Formularvalidierung", "Clientseitige Prüfung von Pflichtfeldern, Streitwert und Datumsangaben vor dem Absenden"),
        ("Kalender-Rendering", "Der Monatskalender wird vollständig im Browser generiert (Funktion renderCalendar)"),
        ("Fristvorschau", "Bei Eingabe des Zustellungsdatums wird das Fristdatum live berechnet und angezeigt"),
        ("Tab-Wechsel", "Umschalten zwischen Dashboard und Kalender ohne Seiten-Reload"),
        ("Task-Hervorhebung", "Bei Navigation aus dem Kalender wird der Ziel-Task per scrollIntoView und Blink-Animation markiert"),
    ]),
    ("⚡ jQuery 3.7", [
        ("8 AJAX-Endpunkte", "create, update, delete, complete für Tasks sowie complete, revert, update, get für Phasen"),
        ("Animationen", "slideToggle für das Ein- und Ausklappen der Phasen-Panels"),
        ("DOM-Manipulation", "Neue Tasks werden dynamisch in die Dashboard-Liste eingefügt"),
        ("Live-Filter", "Die Cases-Tabelle wird per Dropdown nach Phase gefiltert"),
        ("Cascading Dropdown", "Im Kalender-Task-Modal werden Phasen abhängig vom gewählten Fall per AJAX geladen"),
    ]),
])

# ──────────────────────────────────────────────────────────
# FOLIE 17: Vielen Dank
# ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(slide_layouts[19])
for shape in s.shapes:
    if shape.has_text_frame and 'dank' in shape.text.lower():
        shape.text = "Vielen Dank für Ihre Aufmerksamkeit!"
    if shape.is_placeholder and shape.has_text_frame and not shape.text.strip():
        shape.text = "Fragen?"
tb = s.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(7), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "github.com/malala2409/LexCortex---Case-Management"
p.font.size = Pt(12); p.font.color.rgb = MUTED; p.alignment = PP_ALIGN.CENTER

# ── Speichern ────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\n✅ {OUTPUT} — {len(prs.slides)} Folien gespeichert.")
